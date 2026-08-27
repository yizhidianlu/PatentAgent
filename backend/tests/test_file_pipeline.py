"""文件管线专项回归测试（4 号 tester）。

覆盖用户核心痛点「上传文件后功能不可用」的每条失败路径，重点验收：
**失败时给出的必须是人能照做的中文说明，而不是英文堆栈或空提示。**

素材全部由本文件当场生成（pymupdf / python-docx / python-pptx / Pillow），
不依赖任何外部 fixture，也不触发任何 LLM 调用。
"""

from __future__ import annotations

import io
import threading
import zipfile
from pathlib import Path

from conftest import disk_path

import pytest

from app.db import database as db
from app.services import artifacts as artifacts_service
from app.services import convert as convert_service
from app.services import drawings as drawings_service
from app.services import export_docx, export_pdf


@pytest.fixture(scope="module", autouse=True)
def _database():
    """建库跑迁移（本模块不起 FastAPI，直接调 services 层）。"""
    db.init_db()
    yield

# ---------------------------------------------------------------------------
# 公共断言：错误提示的可读性
# ---------------------------------------------------------------------------

# 出现这些片段说明把 Python 堆栈直接丢给了用户
_STACK_MARKERS = (
    "Traceback (most recent call last)",
    "  File \"",
    "~~~~~~",
    "^^^^^^",
)


def assert_readable_chinese_error(msg: str | None, *, must_contain: tuple[str, ...] = ()) -> None:
    """错误提示必须：非空、含中文、无裸堆栈、给出可操作的下一步。"""
    assert msg, "失败路径必须给出错误说明，不能为空"
    assert any("\u4e00" <= ch <= "\u9fff" for ch in msg), f"错误说明必须是中文：{msg!r}"
    for marker in _STACK_MARKERS:
        assert marker not in msg, f"错误说明里混进了 Python 堆栈（{marker!r}）：{msg!r}"
    # 「该怎么办」：至少出现一个祈使性引导词
    assert any(w in msg for w in ("请", "可以", "建议", "解决办法")), f"错误说明没有告诉用户该怎么办：{msg!r}"
    for kw in must_contain:
        assert kw in msg, f"错误说明应提到 {kw!r}：{msg!r}"


# ---------------------------------------------------------------------------
# 素材生成
# ---------------------------------------------------------------------------


def _png_bytes(w: int = 240, h: int = 180, text: str = "FIG") -> bytes:
    from PIL import Image, ImageDraw

    im = Image.new("RGB", (w, h), (210, 210, 210))
    ImageDraw.Draw(im).rectangle([5, 5, w - 5, h - 5], outline=(0, 0, 0), width=3)
    ImageDraw.Draw(im).text((20, 20), text, fill=(0, 0, 0))
    buf = io.BytesIO()
    im.save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture()
def case_dir(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    """把 uploads 指到 tmp_path，返回该案件的上传目录。"""
    from app.config import get_config

    cfg = get_config()
    monkeypatch.setattr(type(cfg), "uploads_dir", property(lambda _self: tmp_path / "uploads"))
    monkeypatch.setattr(type(cfg), "outputs_dir", property(lambda _self: tmp_path / "outputs"))
    d = cfg.uploads_dir / "case_fp"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _convert(case_dir: Path, name: str, payload: bytes) -> convert_service.ConvertResult:
    """把 payload 落成上传件并跑一次真实转换。"""
    path = case_dir / convert_service.sanitize_filename(name)
    path.write_bytes(payload)
    return convert_service.convert_upload("case_fp", path)


# ---------------------------------------------------------------------------
# 1. PDF 转换
# ---------------------------------------------------------------------------


def _pdf(pages: list[tuple[str, ...]], **save_kw) -> bytes:
    import pymupdf

    doc = pymupdf.open()
    for texts in pages:
        page = doc.new_page()
        y = 80
        for t in texts:
            page.insert_text((72, y), t, fontname="china-s", fontsize=12)
            y += 26
    buf = doc.tobytes(**save_kw) if save_kw else doc.tobytes()
    doc.close()
    return buf


def test_pdf_chinese_ok(case_dir: Path) -> None:
    """中文 PDF：正文、页码、图注候选都能抽出来。"""
    data = _pdf([("批任务调度系统技术交底书", "图1 系统总体架构示意图"), ("第二页正文", "表1 参数对照表")])
    res = _convert(case_dir, "中文正常.pdf", data)

    assert res.meta.get("convert_error") is None
    assert res.md_path is not None
    text = res.md_path.read_text(encoding="utf-8")
    assert "批任务调度系统技术交底书" in text
    assert "## 第 1 页" in text and "## 第 2 页" in text
    assert res.meta["pages"] == 2
    captions = {c["text"] for c in res.meta["figure_captions"]}
    assert any("图1" in c for c in captions)
    assert any("表1" in c for c in captions)


def test_pdf_scanned_without_text_layer_is_reported(case_dir: Path) -> None:
    """扫描件（纯图片无文字层）必须明确报错，不能静默产出空 md。

    这是用户「上传 PDF 后功能不可用」最典型的一种：
    以前会返回「上传成功」+ 一份只有分页标题的 md，用户完全不知道发生了什么。
    """
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page(width=595, height=842)
    page.insert_image(pymupdf.Rect(50, 50, 550, 750), stream=_png_bytes(500, 700, "SCAN"))
    data = doc.tobytes()
    doc.close()

    res = _convert(case_dir, "扫描件.pdf", data)

    assert res.md_path is None, "无文字层的 PDF 不应产出可用 md（否则下游会拿到空正文）"
    assert res.meta.get("scanned") is True
    assert res.meta.get("text_chars") == 0
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=("文字层", "OCR"))
    # 页面图片仍然抽出来了，可作为附图素材
    assert res.meta.get("figures"), "扫描件的页面图片应仍被抽取保留"
    assert Path(res.meta["figures"][0]["path"]).is_file()


def test_pdf_blank_page_is_reported(case_dir: Path) -> None:
    """全空白 PDF：既无文字也无图片，同样要说清楚。"""
    import pymupdf

    doc = pymupdf.open()
    doc.new_page()
    data = doc.tobytes()
    doc.close()

    res = _convert(case_dir, "空白.pdf", data)
    assert res.md_path is None
    assert_readable_chinese_error(res.meta.get("convert_error"))


def test_pdf_encrypted_gives_actionable_hint(case_dir: Path) -> None:
    """加密 PDF：提示「需要密码 + 怎么去掉密码」，而不是 `document closed or encrypted`。"""
    import pymupdf

    doc = pymupdf.open()
    doc.new_page().insert_text((72, 90), "secret")
    data = doc.tobytes(
        encryption=pymupdf.PDF_ENCRYPT_AES_256, owner_pw="ownerpw", user_pw="userpw"
    )
    doc.close()

    res = _convert(case_dir, "加密.pdf", data)
    assert res.md_path is None
    assert res.meta.get("encrypted") is True
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=("加密", "密码"))
    assert "encrypted" not in res.meta["convert_error"]


def test_pdf_corrupted_error_is_chinese_and_hides_server_path(case_dir: Path) -> None:
    """损坏 PDF：中文说明，且不回显服务端绝对路径。"""
    res = _convert(case_dir, "损坏.pdf", b"%PDF-1.7\n" + b"\x00\xff garbage " * 50)

    assert res.md_path is None
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=("损坏.pdf",))
    assert "Failed to open file" not in res.meta["convert_error"]
    assert ":\\" not in res.meta["convert_error"], "错误说明不应泄露服务端绝对路径"


def test_pdf_not_really_pdf(case_dir: Path) -> None:
    """扩展名是 .pdf 但内容是纯文本：提示「可能是别的格式改了扩展名」。"""
    res = _convert(case_dir, "伪装.pdf", "这其实是一个文本文件".encode())
    assert res.md_path is None
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=("扩展名",))


def test_pdf_large_document(case_dir: Path) -> None:
    """50+ 页 PDF：全部页都要转出来。"""
    res = _convert(case_dir, "超大.pdf", _pdf([(f"第 {i} 页正文内容",) for i in range(1, 61)]))

    assert res.meta.get("convert_error") is None
    assert res.meta["pages"] == 60
    text = res.md_path.read_text(encoding="utf-8")  # type: ignore[union-attr]
    assert "## 第 1 页" in text and "## 第 60 页" in text
    assert "第 60 页正文内容" in text


def test_pdf_figures_land_in_case_figures_dir(case_dir: Path) -> None:
    """PDF 插图必须落到 uploads/{case}/figures/{主干}/ 并在 md 中以相对路径引用。"""
    import pymupdf

    doc = pymupdf.open()
    for i in (1, 2):
        page = doc.new_page()
        page.insert_text((60, 60), f"图{i} 模块关系", fontname="china-s", fontsize=12)
        page.insert_image(pymupdf.Rect(60, 80, 360, 300), stream=_png_bytes(text=f"IMG{i}"))
    data = doc.tobytes()
    doc.close()

    res = _convert(case_dir, "含图.pdf", data)
    assert res.meta.get("convert_error") is None

    figures = res.meta["figures"]
    assert len(figures) == 2
    expected_dir = case_dir / "figures" / "含图"
    for fig in figures:
        p = Path(fig["path"])
        assert p.parent == expected_dir, f"插图应落在 {expected_dir}，实际 {p.parent}"
        assert p.is_file() and p.stat().st_size > 0
    assert res.meta["media_dir"] == str(expected_dir)

    md_text = res.md_path.read_text(encoding="utf-8")  # type: ignore[union-attr]
    assert "../figures/含图/" in md_text, "md 中应以相对路径引用抽出的插图"
    # 图注候选带页码
    assert {c["page"] for c in res.meta["figure_captions"]} == {1, 2}


def test_pdf_tables_and_formulas_text_preserved(case_dir: Path) -> None:
    """含表格与公式的 PDF：文字内容不丢。"""
    res = _convert(
        case_dir,
        "表格公式.pdf",
        _pdf([("表1 性能对照", "基线 120 800", "本发明 45 2100", "公式：L = Σ w_i · exp(−α t_i) / Z")]),
    )
    assert res.meta.get("convert_error") is None
    text = res.md_path.read_text(encoding="utf-8")  # type: ignore[union-attr]
    for kw in ("表1 性能对照", "本发明", "2100", "Σ"):
        assert kw in text


def test_pdf_english(case_dir: Path) -> None:
    """英文 PDF：Figure N 图注候选同样识别。"""
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 80), "A Method for Scheduling Batch Jobs", fontsize=13)
    page.insert_text((72, 110), "Figure 1 Overall architecture", fontsize=11)
    data = doc.tobytes()
    doc.close()

    res = _convert(case_dir, "english.pdf", data)
    assert res.meta.get("convert_error") is None
    assert any("Figure 1" in c["text"] for c in res.meta["figure_captions"])


# ---------------------------------------------------------------------------
# 2. Word / PPT 转换
# ---------------------------------------------------------------------------


def _docx_bytes(build) -> bytes:
    import docx

    d = docx.Document()
    build(d)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def test_docx_normal_with_headings(case_dir: Path) -> None:
    def build(d):
        d.add_heading("批任务调度方案设计", level=1)
        d.add_paragraph("调度器维护全局队列，按优先级派发任务。")
        d.add_heading("核心机制", level=2)

    res = _convert(case_dir, "正常.docx", _docx_bytes(build))
    assert res.meta.get("convert_error") is None
    text = res.md_path.read_text(encoding="utf-8")  # type: ignore[union-attr]
    assert "# 批任务调度方案设计" in text
    assert "调度器维护全局队列" in text


def test_docx_with_image_extracts_media(case_dir: Path, tmp_path: Path) -> None:
    from docx.shared import Inches

    img = tmp_path / "embed.png"
    img.write_bytes(_png_bytes(text="DOCX"))

    def build(d):
        d.add_paragraph("含图片的文档")
        d.add_picture(str(img), width=Inches(2))

    res = _convert(case_dir, "带图.docx", _docx_bytes(build))
    assert res.meta.get("convert_error") is None
    assert res.meta.get("figure_count", 0) >= 1
    media = disk_path(res.meta["media_dir"])
    assert media.is_dir() and any(media.iterdir())


def test_docx_with_table(case_dir: Path) -> None:
    def build(d):
        d.add_paragraph("含表格的文档")
        t = d.add_table(rows=2, cols=2)
        t.cell(0, 0).text = "方法"
        t.cell(0, 1).text = "时延"
        t.cell(1, 0).text = "本发明"
        t.cell(1, 1).text = "45"

    res = _convert(case_dir, "带表.docx", _docx_bytes(build))
    assert res.meta.get("convert_error") is None
    text = res.md_path.read_text(encoding="utf-8")  # type: ignore[union-attr]
    assert "本发明" in text and "45" in text


def test_docx_broken_error_is_actionable_chinese(case_dir: Path) -> None:
    """坏 docx：以前直接把 zipfile.BadZipFile 的 Python 堆栈丢给用户。"""
    res = _convert(case_dir, "坏文件.docx", b"PK\x03\x04 not a real docx")

    assert res.md_path is None
    err = res.meta.get("convert_error")
    assert_readable_chinese_error(err, must_contain=("坏文件.docx", "另存为", ".docx"))
    assert "Traceback" not in err and "site-packages" not in err  # type: ignore[operator]


def test_docx_empty_body_is_reported(case_dir: Path) -> None:
    """能打开但一个字都没有的 docx：不能算「转换成功」。"""
    res = _convert(case_dir, "空内容.docx", _docx_bytes(lambda d: None))
    assert res.md_path is None, "只剩转换器自带注释的 md 等同于没转出东西"
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=("空内容.docx",))


def test_legacy_doc_format_tells_user_to_save_as_docx(case_dir: Path) -> None:
    """旧版 .doc：明确「不支持 + 另存为 .docx」，而不是含糊的「仅存原件」。"""
    ole2 = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 512
    res = _convert(case_dir, "老格式.doc", ole2)

    assert res.md_path is None
    assert res.meta["category"] == "legacy_office"
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=("另存为", ".docx"))


def test_legacy_ppt_format_tells_user_to_save_as_pptx(case_dir: Path) -> None:
    ole2 = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 512
    res = _convert(case_dir, "老格式.ppt", ole2)

    assert res.md_path is None
    assert res.meta["category"] == "legacy_office"
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=("另存为", ".pptx"))


def _pptx_bytes(build) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    build(prs)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_normal(case_dir: Path) -> None:
    def build(prs):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "批任务调度方案评审"
        s.placeholders[1].text = "异构节点资源错配问题"

    res = _convert(case_dir, "正常.pptx", _pptx_bytes(build))
    assert res.meta.get("convert_error") is None
    text = res.md_path.read_text(encoding="utf-8")  # type: ignore[union-attr]
    assert "## 第 1 页" in text
    assert "批任务调度方案评审" in text and "异构节点资源错配问题" in text


def test_pptx_with_picture(case_dir: Path, tmp_path: Path) -> None:
    from pptx.util import Inches

    img = tmp_path / "slide.png"
    img.write_bytes(_png_bytes(text="PPT"))

    def build(prs):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "含图表的幻灯片"
        s.shapes.add_picture(str(img), Inches(1), Inches(2), width=Inches(3))

    res = _convert(case_dir, "带图.pptx", _pptx_bytes(build))
    assert res.meta.get("convert_error") is None
    assert res.meta.get("figure_count", 0) >= 1


def test_pptx_broken_error_is_actionable_chinese(case_dir: Path) -> None:
    res = _convert(case_dir, "坏演示.pptx", b"PK\x03\x04 broken")
    assert res.md_path is None
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=("坏演示.pptx", ".pptx"))


# ---------------------------------------------------------------------------
# 3. 图片与其它格式
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    ("name", "fmt"),
    [("图.png", "PNG"), ("图.jpg", "JPEG"), ("图.webp", "WEBP")],
)
def test_images_direct_store_with_metadata(case_dir: Path, name: str, fmt: str) -> None:
    """位图直存：无 md、无错误，并记录尺寸/格式供后续附图使用。"""
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (120, 90), (20, 90, 160)).save(buf, format=fmt)
    res = _convert(case_dir, name, buf.getvalue())

    assert res.md_path is None
    assert res.meta.get("convert_error") is None
    assert res.meta["category"] == "image"
    assert res.meta["image_size"] == [120, 90]
    assert res.meta["image_format"] == fmt


def test_svg_direct_store(case_dir: Path) -> None:
    svg = b'<svg xmlns="http://www.w3.org/2000/svg" width="100" height="50"><rect/></svg>'
    res = _convert(case_dir, "矢量.svg", svg)
    assert res.meta.get("convert_error") is None
    assert res.meta["category"] == "image"


def test_svg_broken_is_reported(case_dir: Path) -> None:
    res = _convert(case_dir, "假矢量.svg", b"this is not svg at all")
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=("SVG",))


def test_huge_image_ok(case_dir: Path) -> None:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (4000, 3000), (0, 0, 0)).save(buf, format="PNG")
    res = _convert(case_dir, "超大图.png", buf.getvalue())
    assert res.meta.get("convert_error") is None
    assert res.meta["image_size"] == [4000, 3000]


def test_corrupt_image_is_reported(case_dir: Path) -> None:
    """损坏图片：上传当场就说清楚，而不是等嵌进 Word 时才炸。"""
    res = _convert(case_dir, "坏图.png", b"\x89PNG\r\n\x1a\n" + b"garbage" * 20)
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=("坏图.png",))
    assert ":\\" not in res.meta["convert_error"], "错误说明不应泄露服务端绝对路径"


@pytest.mark.parametrize(
    ("name", "payload"),
    [
        ("说明.md", "# 标题\n\n正文段落\n"),
        ("笔记.txt", "纯文本第一行\n第二行\n"),
        ("代码.py", "def f(x):\n    return x * 2\n"),
    ],
)
def test_text_like_files_direct_store(case_dir: Path, name: str, payload: str) -> None:
    """文本/代码/md 直存：md_path 指向原件本身。"""
    res = _convert(case_dir, name, payload.encode("utf-8"))
    assert res.meta.get("convert_error") is None
    assert res.meta["category"] == "text"
    assert res.md_path is not None and res.md_path.name == name


def test_empty_file_is_reported(case_dir: Path) -> None:
    """0 字节文件：明确告知，而不是产出一份空 md。"""
    res = _convert(case_dir, "空文件.txt", b"")
    assert res.md_path is None
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=("0 字节",))


def test_unknown_extension_is_reported(case_dir: Path) -> None:
    """不支持的扩展名：说明「原样保存 + 可以转成什么格式」。"""
    res = _convert(case_dir, "数据.xyz", b"binary payload")
    assert res.md_path is None
    assert res.meta["category"] == "other"
    assert_readable_chinese_error(res.meta.get("convert_error"), must_contain=(".xyz", ".docx"))


# ---------------------------------------------------------------------------
# 4. 文件名边界
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    ("raw", "expected"),
    [
        ("中文文件名.md", "中文文件名.md"),
        ("有 空 格 的 文件.md", "有 空 格 的 文件.md"),
        ("emoji🚀名字✨.md", "emoji🚀名字✨.md"),
        ("../../escape.md", "escape.md"),          # 路径穿越只留 basename
        ("C:\\Windows\\evil.md", "evil.md"),
        ("tab\there.md", "tabhere.md"),            # 控制字符剔除
        ("dot.in.middle.name.md", "dot.in.middle.name.md"),
        ("....md", "....md"),                       # 前导点保留，扩展名不能被削掉
        ("...", "unnamed"),
        ("", "unnamed"),
    ],
)
def test_sanitize_filename(raw: str, expected: str) -> None:
    assert convert_service.sanitize_filename(raw) == expected


@pytest.mark.parametrize("stem_len", [120, 200, 400])
def test_long_filename_is_truncated_and_uploadable(case_dir: Path, stem_len: int) -> None:
    """超长文件名（期刊 PDF 的常态）必须能落盘。

    以前不截断，Windows 上整条路径破 MAX_PATH(260)，
    `stored_path.write_bytes()` 抛 FileNotFoundError 把整个上传请求打成 500。
    """
    raw = "A" * stem_len + ".pdf"
    safe = convert_service.sanitize_filename(raw)
    assert len(safe) <= convert_service.MAX_FILENAME_LEN
    assert safe.endswith(".pdf"), "截断后必须保留扩展名，否则会被误判为不支持的格式"

    path = convert_service.allocate_upload_path(case_dir, raw)
    path.write_bytes(_pdf([("长文件名测试",)]))       # 不抛 FileNotFoundError 即通过
    res = convert_service.convert_upload("case_fp", path)
    assert res.meta.get("convert_error") is None
    assert res.md_path is not None


def test_long_chinese_filename_is_truncated(case_dir: Path) -> None:
    safe = convert_service.sanitize_filename("中" * 150 + ".docx")
    assert len(safe) <= convert_service.MAX_FILENAME_LEN
    assert safe.endswith(".docx")


def test_duplicate_names_get_suffixed(case_dir: Path) -> None:
    """同名上传：自动 _1/_2，绝不覆盖。"""
    names = []
    for _ in range(3):
        p = convert_service.allocate_upload_path(case_dir, "dup.md")
        p.write_text("x", encoding="utf-8")
        names.append(p.name)
    assert names == ["dup.md", "dup_1.md", "dup_2.md"]


# ---------------------------------------------------------------------------
# 5. md → docx 导出
# ---------------------------------------------------------------------------

_RICH_MD = """**案件名称**：一种批任务调度方法

# 一级标题

正文中文段落，含**加粗**与`行内代码`。

## 二级标题

### 三级标题

行内公式 $E = mc^2$ 与行间公式：

$$
L = \\sum_{i=1}^{n} w_i \\exp(-\\alpha t_i) / Z
$$

| 方法 | 时延 | 吞吐 |
| --- | --- | --- |
| 基线 | 120 | 800 |
| 本发明 | 45 | 2100 |

```python
def schedule(jobs):
    return sorted(jobs, key=lambda j: j.priority)
```

- 列表项一
- 列表项二
"""


def _docx_xml(path: Path) -> tuple[str, list[str]]:
    with zipfile.ZipFile(path) as z:
        return z.read("word/document.xml").decode("utf-8", "replace"), z.namelist()


def test_export_md_to_docx_rich_content(tmp_path: Path) -> None:
    """中文 / 多级标题 / 表格 / 代码块 / LaTeX 公式：公式必须是可编辑 OMML，不是图片。"""
    md = tmp_path / "rich.md"
    md.write_text(_RICH_MD, encoding="utf-8")
    out = tmp_path / "rich.docx"

    stats = export_docx.export_md_to_docx_sync(md, out)
    assert out.is_file() and out.stat().st_size > 0

    xml, _ = _docx_xml(out)
    assert "<m:oMath" in xml, "LaTeX 公式应转成可编辑的 OMML"
    assert stats.get("math", {}).get("omml", 0) >= 1
    assert stats["math"]["png"] == 0, "公式不应降级成图片"

    from docx import Document

    texts = [p.text for p in Document(str(out)).paragraphs]
    joined = "\n".join(texts)
    assert "一级标题" in joined and "三级标题" in joined
    assert "def schedule(jobs):" in joined
    tables = Document(str(out)).tables
    assert tables and any("本发明" in c.text for c in tables[0]._cells)


def test_export_md_to_docx_long_document(tmp_path: Path) -> None:
    """长文档（120 段）：段落不丢。"""
    body = "\n\n".join(f"## 第 {i} 节\n\n这是第 {i} 段正文内容。" for i in range(1, 121))
    md = tmp_path / "long.md"
    md.write_text("**案件名称**：长文档\n\n" + body, encoding="utf-8")
    out = tmp_path / "long.docx"

    export_docx.export_md_to_docx_sync(md, out)
    from docx import Document

    texts = "\n".join(p.text for p in Document(str(out)).paragraphs)
    assert "第 1 节" in texts and "第 120 节" in texts


def test_export_md_to_docx_mermaid_renders_png(tmp_path: Path) -> None:
    """mermaid 围栏渲染成 PNG 并嵌入 Word；浏览器不可用则 skip。"""
    md = tmp_path / "m.md"
    md.write_text(
        "# 图\n\n```mermaid\nflowchart TD\n  A[接收任务] --> B[派发]\n```\n", encoding="utf-8"
    )
    out = tmp_path / "m.docx"
    stats = export_docx.export_md_to_docx_sync(md, out)
    if stats.get("mermaid_error") or stats.get("mermaid_ok", 0) < 1:
        pytest.skip(f"本机 mermaid 渲染不可用：{stats}")

    _, names = _docx_xml(out)
    assert any(n.startswith("word/media/") for n in names), "渲染出的 mermaid PNG 应嵌进 docx"


def test_export_md_to_docx_bad_mermaid_degrades(tmp_path: Path) -> None:
    """语法错误的 mermaid 单块降级为源码，不能让整份导出失败。"""
    md = tmp_path / "bad.md"
    md.write_text("# 坏图\n\n```mermaid\n!!! not valid @@@\n```\n\n正文仍在。\n", encoding="utf-8")
    out = tmp_path / "bad.docx"

    stats = export_docx.export_md_to_docx_sync(md, out)
    assert out.is_file()
    if not stats.get("mermaid_error"):
        assert stats["mermaid_fail"] >= 1
    from docx import Document

    assert any("正文仍在" in p.text for p in Document(str(out)).paragraphs)


def test_export_md_to_docx_missing_input_error(tmp_path: Path) -> None:
    with pytest.raises(export_docx.DocxExportError) as ei:
        export_docx.export_md_to_docx_sync(tmp_path / "nope.md", tmp_path / "nope.docx")
    assert_readable_chinese_error(str(ei.value), must_contain=("nope.md",))


def test_export_md_to_docx_empty_input_error(tmp_path: Path) -> None:
    """空 md 不应静默产出一份空白 Word。"""
    md = tmp_path / "empty.md"
    md.write_text("   \n\n", encoding="utf-8")
    with pytest.raises(export_docx.DocxExportError) as ei:
        export_docx.export_md_to_docx_sync(md, tmp_path / "empty.docx")
    assert_readable_chinese_error(str(ei.value), must_contain=("空白",))


# ---------------------------------------------------------------------------
# 6. docx → pdf 导出与降级链
# ---------------------------------------------------------------------------


def test_pdf_probe_structure() -> None:
    report = export_pdf.probe()
    assert set(report) >= {"word", "soffice", "order"}
    assert isinstance(report["order"], list)
    assert all(e in ("word", "soffice") for e in report["order"])


async def test_docx_to_pdf_missing_input(tmp_path: Path) -> None:
    with pytest.raises(export_pdf.PdfExportError) as ei:
        await export_pdf.docx_to_pdf(tmp_path / "nope.docx", tmp_path / "nope.pdf")
    assert_readable_chinese_error(str(ei.value), must_contain=("nope.docx",))


async def test_docx_to_pdf_unsupported_engine_setting(tmp_path: Path) -> None:
    """settings.general.pdf_engine=pillow：以前报「全部引擎失败：。」（原因为空）。"""
    docx = tmp_path / "x.docx"
    docx.write_bytes(b"placeholder")
    with pytest.raises(export_pdf.PdfExportError) as ei:
        await export_pdf.docx_to_pdf(docx, tmp_path / "x.pdf", engine="pillow")
    msg = str(ei.value)
    assert_readable_chinese_error(msg, must_contain=("pillow", "设置"))
    assert "：。" not in msg, "错误信息里不能出现空的原因列表"


async def test_docx_to_pdf_all_engines_unavailable(tmp_path: Path, monkeypatch) -> None:
    """Word 与 soffice 都不可用时：给出安装指引与「先下载 docx」的替代方案。"""
    monkeypatch.setattr(export_pdf, "_word_available", lambda: False)
    monkeypatch.setattr(export_pdf, "_soffice_path", lambda: None)
    docx = tmp_path / "x.docx"
    docx.write_bytes(b"placeholder")

    with pytest.raises(export_pdf.PdfExportError) as ei:
        await export_pdf.docx_to_pdf(docx, tmp_path / "x.pdf")
    assert_readable_chinese_error(str(ei.value), must_contain=("Word", "LibreOffice", "docx"))


async def test_docx_to_pdf_falls_back_to_soffice_when_word_fails(tmp_path: Path, monkeypatch) -> None:
    """降级链 word → soffice：Word 报错后必须真的去试 soffice。"""
    calls: list[str] = []

    def fake_word(docx_path: Path, pdf_path: Path) -> None:
        calls.append("word")
        raise RuntimeError("Word 模拟不可用")

    def fake_soffice(docx_path: Path, pdf_path: Path, soffice: str) -> None:
        calls.append("soffice")
        pdf_path.parent.mkdir(parents=True, exist_ok=True)
        pdf_path.write_bytes(b"%PDF-1.4 fake")

    monkeypatch.setattr(export_pdf, "_word_available", lambda: True)
    monkeypatch.setattr(export_pdf, "_docx_to_pdf_word_sync", fake_word)
    monkeypatch.setattr(export_pdf, "_soffice_path", lambda: "soffice")
    monkeypatch.setattr(export_pdf, "_docx_to_pdf_soffice_sync", fake_soffice)

    docx = tmp_path / "x.docx"
    docx.write_bytes(b"placeholder")
    engine = await export_pdf.docx_to_pdf(docx, tmp_path / "x.pdf")

    assert engine == "soffice"
    assert calls == ["word", "soffice"]


async def test_docx_to_pdf_com_error_is_readable(tmp_path: Path, monkeypatch) -> None:
    """Word COM 异常必须压成一句人话，而不是把整个 COM 元组和 chm 文件名丢给用户。"""

    class FakeComError(Exception):
        pass

    exc = FakeComError(
        -2147352567,
        "发生意外。",
        (0, "Microsoft Word", "Word 在试图打开文件时遇到错误。\n请检查文档权限。", "wdmain11.chm", 24601, -1),
        None,
    )

    def boom(docx_path: Path, pdf_path: Path) -> None:
        raise exc

    monkeypatch.setattr(export_pdf, "_word_available", lambda: True)
    monkeypatch.setattr(export_pdf, "_docx_to_pdf_word_sync", boom)
    monkeypatch.setattr(export_pdf, "_soffice_path", lambda: None)

    docx = tmp_path / "x.docx"
    docx.write_bytes(b"placeholder")
    with pytest.raises(export_pdf.PdfExportError) as ei:
        await export_pdf.docx_to_pdf(docx, tmp_path / "x.pdf")

    msg = str(ei.value)
    assert "Word 在试图打开文件时遇到错误。 请检查文档权限。" in msg
    assert "wdmain11.chm" not in msg, "帮助文件名对用户毫无意义"
    assert "-2147352567" not in msg, "COM HRESULT 对用户毫无意义"


def test_humanize_com_error_shapes() -> None:
    """_humanize_com_error 对 excepinfo / args 两种形态都要能提取描述。"""

    class E(Exception):
        pass

    e1 = E(-1, "发生意外。", (0, "Microsoft Word", "真正的说明", "x.chm", 1, 2), None)
    assert export_pdf._humanize_com_error(e1) == "Microsoft Word：真正的说明"

    e2 = E("只有一句话")
    e2.excepinfo = (0, "Word", "来自 excepinfo 的说明", "x.chm", 1, 2)  # type: ignore[attr-defined]
    assert export_pdf._humanize_com_error(e2) == "Word：来自 excepinfo 的说明"

    e3 = E("纯字符串异常")
    assert export_pdf._humanize_com_error(e3) == "纯字符串异常"


# ---------------------------------------------------------------------------
# 7. 附图生成（零依赖脚本，真跑）
# ---------------------------------------------------------------------------

_DRAWING_CONTENT = {
    "drawings": [
        (
            "图1：一种超声图像病灶分割方法流程图，包含步骤S101，对超声图像进行超像素划分；"
            "S102，构建超像素邻接图；S103，图卷积网络逐层聚合邻域特征；S104，生成病灶分割掩膜。"
        ),
        (
            "图2：一种超声图像病灶分割系统结构示意图，包含超像素划分模块、邻接关系构建模块、"
            "特征聚合分类模块、掩膜生成模块。"
        ),
        "图3：数据流示意图，候选动作数据流经特征提取单元、策略参数更新单元、动作输出单元。",
    ],
    "description": {"drawing_description": ["图1为流程图。", "图2为结构示意图。", "图3为数据流图。"]},
}


@pytest.fixture()
def outputs_case(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> str:
    from app.config import get_config

    cfg = get_config()
    monkeypatch.setattr(type(cfg), "outputs_dir", property(lambda _self: tmp_path / "outputs"))
    return "case_dw"


def test_drawings_three_layouts(outputs_case: str) -> None:
    """method_flow / system_block / data_flow 三种版式：SVG 合法、PNG 可打开、
    画布占比校验通过、图内不含图号。"""
    import xml.etree.ElementTree as ET

    from PIL import Image

    updated = drawings_service.run_generator_sync(outputs_case, _DRAWING_CONTENT)
    assets = updated.get("drawing_assets") or []
    assert len(assets) == 3
    assert drawings_service.failed_figures(updated) == []

    for asset in assets:
        v = asset["validation"]
        assert v["passes"] is True
        # 画布占比校验真的生效（不低于 0.80 才 passes）
        assert v["content_width_ratio"] >= 0.79 and v["content_height_ratio"] >= 0.79
        assert v["internal_title"] is False

    files = drawings_service.asset_files(outputs_case, updated)
    assert len(files) == 3
    for entry in files:
        svg = entry["svg_path"]
        root = ET.parse(svg).getroot()
        assert root.tag == "{http://www.w3.org/2000/svg}svg"
        labels = [(t.text or "").strip() for t in root.iter("{http://www.w3.org/2000/svg}text")]
        assert labels, "SVG 里应有节点文字"
        # 图内不得出现图号 / 图题
        assert not [t for t in labels if t[:1] == "图" and t[1:2].isdigit()]

        png = entry["png_path"]
        with Image.open(png) as im:
            im.verify()
        with Image.open(png) as im:
            assert im.size[0] > 0 and im.size[1] > 0


def test_drawings_unparsable_spec_error_is_chinese(outputs_case: str) -> None:
    """规格里识别不出步骤/模块时，报错要说明是哪张图、该怎么改写规格。"""
    with pytest.raises(drawings_service.DrawingGenerationError) as ei:
        drawings_service.run_generator_sync(
            outputs_case, {"drawings": ["图1：一种示意图，其内容无法解析出步骤或模块。"]}
        )
    exc = ei.value
    assert exc.figure_no == 1
    assert_readable_chinese_error(str(exc), must_contain=("图1", "步骤"))
    assert "Traceback" not in str(exc) and ".py" not in str(exc)


def test_drawings_empty_list_is_noop(outputs_case: str) -> None:
    updated = drawings_service.run_generator_sync(outputs_case, {"drawings": []})
    assert (updated.get("drawing_assets") or []) == []


def test_drawings_degrade_figure(outputs_case: str) -> None:
    """降级：移出 drawings/assets/附图说明，补 image_model_prompt 与中文 gaps 文案。"""
    content = {
        "drawings": ["图1：流程图，包含步骤S101，甲；S102，乙。", "图2：结构示意图，包含A模块、B模块。"],
        "drawing_assets": [{"figure_no": 2, "spec": "图2：结构示意图，包含A模块、B模块。"}],
        "description": {"drawing_description": ["图1为流程图。", "图2为结构示意图。"]},
    }
    gap = drawings_service.degrade_figure(content, 2, "规格不足以成图")

    assert len(content["drawings"]) == 1 and not content["drawing_assets"]
    assert content["description"]["drawing_description"] == ["图1为流程图。"]
    prompts = content["image_model_prompts"]
    assert [p["figure_no"] for p in prompts] == [2]
    assert "黑白线条" in prompts[0]["prompt"]
    assert "图2" in gap and gap in content["gaps"]


# ---------------------------------------------------------------------------
# 8. artifacts 版本化落盘
# ---------------------------------------------------------------------------


@pytest.fixture()
def artifact_case(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> str:
    """建一个真实 case 行并把 outputs 指到 tmp_path。"""
    from app.config import get_config

    cfg = get_config()
    monkeypatch.setattr(type(cfg), "outputs_dir", property(lambda _self: tmp_path / "outputs"))
    case_id = f"case_art_{tmp_path.name}"[:26]
    db.execute(
        "INSERT INTO cases(id, module, title, status, created_at, updated_at) VALUES (?,?,?,?,?,?)",
        (case_id, "disclosure", "落盘测试", "draft", db.now_str(), db.now_str()),
    )
    return case_id


def test_artifact_concurrent_saves_never_overwrite(artifact_case: str) -> None:
    """并发落盘绝不互相覆盖（§2「只增不改、禁止覆盖」）。

    以前 _allocate_output_path 先 exists() 再 write，存在 TOCTOU 竞态：
    12 个并发落盘在 DB 里留下 12 条版本记录，磁盘上却只剩 4 个文件，
    8 份内容被静默覆盖丢失。
    """
    n = 12
    saved: list = []
    errors: list[str] = []
    barrier = threading.Barrier(n)

    def worker(i: int) -> None:
        barrier.wait()
        try:
            saved.append(
                artifacts_service.save_artifact_sync(
                    artifact_case, "disclosure_md", f"**案件名称**：并发\n\nBODY-{i:03d}", "md"
                )
            )
        except Exception as exc:  # noqa: BLE001
            errors.append(f"{type(exc).__name__}: {exc}")

    threads = [threading.Thread(target=worker, args=(i,)) for i in range(n)]
    for t in threads:
        t.start()
    for t in threads:
        t.join()

    assert not errors, errors
    assert len(saved) == n
    assert len({a.version for a in saved}) == n, "版本号必须唯一"
    assert len({a.filename for a in saved}) == n, "文件名必须唯一"

    bodies = [disk_path(a.stored_path).read_text(encoding="utf-8").strip().splitlines()[-1] for a in saved]
    assert len(set(bodies)) == n, "每条版本记录都必须指向自己那份内容，不能被别人覆盖"

    from app.config import get_config

    on_disk = list((get_config().outputs_dir / artifact_case).iterdir())
    assert len(on_disk) == n


def test_artifact_sequential_saves_increment_version(artifact_case: str) -> None:
    a1 = artifacts_service.save_artifact_sync(artifact_case, "revision_log_md", "第一版", "md", title="顺序测试")
    a2 = artifacts_service.save_artifact_sync(artifact_case, "revision_log_md", "第二版", "md", title="顺序测试")
    assert a2.version == a1.version + 1
    assert a1.filename != a2.filename
    assert disk_path(a1.stored_path).read_text(encoding="utf-8") == "第一版"
    assert disk_path(a2.stored_path).read_text(encoding="utf-8") == "第二版"


@pytest.mark.parametrize(
    ("md_text", "title", "expected"),
    [
        ("**案件名称**：一种批任务调度方法\n\n正文", None, "一种批任务调度方法"),
        ("案件名称：无粗体也可以\n正文", None, "无粗体也可以"),
        ("**案件名称**：XXX\n\n正文", "备用标题", "备用标题"),
        ("**案件名称**：待填写\n\n正文", None, "未命名案件"),
        ("**案件名称**：×××\n\n正文", None, "未命名案件"),
        (None, 'A/B\\C:D*E?F"G<H>I|J', "ABCDEFGHIJ"),
        ("**案件名称**：**《一种方法》**\n\n正文", None, "《一种方法》"),
        (None, None, "未命名案件"),
    ],
)
def test_normalize_case_name_cases(md_text: str | None, title: str | None, expected: str) -> None:
    assert artifacts_service.normalize_case_name(md_text, title) == expected


def test_case_name_line_does_not_swallow_next_paragraph() -> None:
    """「**案件名称**：」后面为空时，不能把下一段正文当成案件名。

    原正则 `[:：]\\s*(.+?)` 里的 \\s* 会跨行，导致文件名变成正文首句。
    """
    assert artifacts_service.normalize_case_name("**案件名称**：\n\n这是正文第一段", None) == "未命名案件"


def test_normalize_case_name_length_cap() -> None:
    assert len(artifacts_service.normalize_case_name(None, "长" * 300)) == artifacts_service.MAX_NAME_LEN


def test_artifact_long_and_illegal_case_name_lands_on_disk(artifact_case: str) -> None:
    """案件名超长 / 含 Windows 非法字符时仍能落盘。"""
    a = artifacts_service.save_artifact_sync(
        artifact_case, "disclosure_md", "**案件名称**：" + "超" * 300 + "\n\n正文", "md"
    )
    assert disk_path(a.stored_path).is_file()
    # {≤80 字案件名}_{14 位时间戳}.md
    assert len(Path(a.filename).stem) <= artifacts_service.MAX_NAME_LEN + 15

    b = artifacts_service.save_artifact_sync(
        artifact_case, "disclosure_md", '**案件名称**：A/B:C*D?E"F<G>H|I\n\n正文', "md"
    )
    assert disk_path(b.stored_path).is_file()
    assert not set(Path(b.filename).stem) & set('\\/:*?"<>|')


def test_artifact_unwritable_dir_error_is_readable(artifact_case: str, monkeypatch) -> None:
    """输出目录不可写时给出中文说明，而不是裸 WinError。"""
    from app.config import get_config

    cfg = get_config()
    monkeypatch.setattr(type(cfg), "outputs_dir", property(lambda _self: Path("Z:/no/such/place")))
    with pytest.raises(OSError) as ei:
        artifacts_service.save_artifact_sync(artifact_case, "disclosure_md", "正文", "md")
    assert_readable_chinese_error(str(ei.value), must_contain=("写入",))


def test_strip_timestamp() -> None:
    assert artifacts_service.strip_timestamp("一种方法_20250101120000.md") == "一种方法"
    assert artifacts_service.strip_timestamp("没有时间戳.md") == "没有时间戳"


# ---------------------------------------------------------------------------
# 9. mermaid 渲染错误提示
# ---------------------------------------------------------------------------


def test_mermaid_empty_code_error() -> None:
    from app.services import mermaid as mermaid_service

    with pytest.raises(mermaid_service.MermaidRenderError) as ei:
        mermaid_service.render_mermaid_png_sync("   ")
    assert "空" in str(ei.value)


@pytest.mark.parametrize(
    ("stderr", "expect_hint"),
    [
        (
            (
                "[mermaid_render] 第 1 个 mermaid 围栏生图失败（已保留源码）："
                "No diagram type detected matching given configuration for text: A --> B\n"
                r"已写入 C:\Users\me\data\tmp\mermaid_01ABC\out.md（mermaid：0 处已转为 PNG，1 处生图失败）"
            ),
            "图表类型声明",
        ),
        (
            "[mermaid_render] 生图失败：Parse error on line 3: expecting SEMI",
            "语法有误",
        ),
        (
            "[mermaid_render] 无法启动出图（将保留 mermaid 源码）：Executable doesn't exist at ...",
            "浏览器",
        ),
    ],
)
def test_mermaid_error_is_translated_and_path_free(stderr: str, expect_hint: str) -> None:
    """mermaid 报错要翻成中文指引，且不能把服务端临时目录回显给用户。"""
    from app.services.mermaid import _extract_error

    msg = _extract_error(stderr)
    assert expect_hint in msg
    assert "已写入" not in msg
    assert "mermaid_01ABC" not in msg, "错误说明不应回显服务端临时目录"
    assert "out.md" not in msg


def test_mermaid_error_never_empty() -> None:
    from app.services.mermaid import _extract_error

    assert _extract_error("").strip()
    assert _extract_error("   \n\n  ").strip()


# ---------------------------------------------------------------------------
# 10. 端到端：上传 API 不因文件本身的问题而 500
# ---------------------------------------------------------------------------


@pytest.fixture(scope="module")
def api_client(client):
    """复用 conftest 的已登录普通用户（上传接口现在要求认证）。"""
    return client


@pytest.fixture(scope="module")
def api_case(api_client) -> str:
    resp = api_client.post("/api/v1/cases", json={"module": "disclosure", "title": "文件管线专项"})
    assert resp.status_code == 201
    return resp.json()["id"]


def _api_upload(client, case_id: str, filename: str, payload: bytes, mime: str) -> dict:
    resp = client.post(
        f"/api/v1/cases/{case_id}/files",
        files=[("files", (filename, io.BytesIO(payload), mime))],
    )
    assert resp.status_code == 201, resp.text
    return resp.json()[0]


def test_api_upload_very_long_filename(api_client, api_case: str) -> None:
    """超长文件名（期刊 PDF 常态）不能把整个上传请求打成 500。"""
    name = "A_very_long_paper_title_" * 12 + ".pdf"      # ~300 字符
    assert len(name) > 260
    item = _api_upload(api_client, api_case, name, _pdf([("长文件名测试正文",)]), "application/pdf")

    assert item["convert_error"] is None
    stored = disk_path(item["file"]["stored_path"])
    assert stored.is_file()
    assert len(stored.name) <= convert_service.MAX_FILENAME_LEN


def test_api_upload_batch_survives_one_bad_file(api_client, api_case: str) -> None:
    """一批文件里混入坏文件：其余文件照常转换，坏文件带中文 convert_error 返回。"""
    good = _pdf([("正常文件正文",)])
    resp = api_client.post(
        f"/api/v1/cases/{api_case}/files",
        files=[
            ("files", ("好文件.pdf", io.BytesIO(good), "application/pdf")),
            ("files", ("坏文件.docx", io.BytesIO(b"not a docx"), "application/octet-stream")),
            ("files", ("空文件.txt", io.BytesIO(b""), "text/plain")),
        ],
    )
    assert resp.status_code == 201, resp.text
    items = resp.json()
    assert len(items) == 3

    ok, bad, empty = items
    assert ok["convert_error"] is None and ok["md_preview"]
    assert_readable_chinese_error(bad["convert_error"])
    assert_readable_chinese_error(empty["convert_error"], must_contain=("0 字节",))
    # 原件一律落盘，用户还能下载回去
    for it in items:
        assert disk_path(it["file"]["stored_path"]).is_file()


def test_api_upload_scanned_pdf_surfaces_reason(api_client, api_case: str) -> None:
    """扫描件经 API 上传：convert_error 直达前端，content 接口给出 404 而非空文本。"""
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page(width=595, height=842)
    page.insert_image(pymupdf.Rect(50, 50, 550, 750), stream=_png_bytes(400, 560, "SCAN"))
    data = doc.tobytes()
    doc.close()

    item = _api_upload(api_client, api_case, "扫描件.pdf", data, "application/pdf")
    assert_readable_chinese_error(item["convert_error"], must_contain=("文字层",))
    assert item["file"]["md_path"] is None
    assert item["md_preview"] is None
    assert api_client.get(f"/api/v1/files/{item['file']['id']}/content").status_code == 404
    # 原件仍可下载
    assert api_client.get(f"/api/v1/files/{item['file']['id']}/download").status_code == 200
